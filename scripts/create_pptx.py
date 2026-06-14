#!/usr/bin/env python3
"""Generate RetroMindAI final presentation .pptx for 16 June 2026 jury."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCREENSHOTS = "/Users/ekeministephen/PycharmProjects/RetroMindAI/docs/presentation/screenshots"

# ── Brand colours ──────────────────────────────────────────────
TEAL      = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6B)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1E, 0x20, 0x2B)
GREY      = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG  = RGBColor(0xF8, 0xFA, 0xFC)
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN     = RGBColor(0x10, 0xB9, 0x81)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_slide(slide, title, bullets, subtitle=None, title_color=TEAL):
    # Title bar
    add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
                 title, size=32, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
                     subtitle, size=14, color=RGBColor(0xCC,0xDD,0xDD))
    # Accent line
    add_rect(slide, 0, Inches(1.1), Inches(0.15), H - Inches(1.1), AMBER)
    # Bullets
    y = Inches(1.6)
    for b in bullets:
        txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        y += Inches(0.45)

def add_image_slide(slide, title, img_path, caption=None):
    add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
                 title, size=32, bold=True, color=WHITE)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
                                 Inches(1.5), Inches(1.6),
                                 Inches(10.3), Inches(5.0))
    if caption:
        add_text_box(slide, Inches(1.5), Inches(6.7), Inches(10), Inches(0.5),
                     caption, size=12, color=GREY, align=PP_ALIGN.CENTER)

def add_footer(slide, slide_num):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.3),
                 f"RetroMindAI  |  Panel 2  |  16 June 2026",
                 size=9, color=GREY)
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
                 str(slide_num), size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(slide, 0, 0, W, H, DARK_BG)
add_rect(slide, 0, Inches(2.0), W, Inches(3.5), TEAL)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2),
             "RetroMindAI", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
             "AI-Powered EV Retrofit Assessment Platform",
             size=28, color=RGBColor(0xCC,0xEE,0xEE), align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
             "Final Presentation  |  Panel 2  |  16 June 2026",
             size=16, color=AMBER, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.1), Inches(11.3), Inches(0.4),
             "Virtual — meet.google.com/xjc-jogb-trd",
             size=12, color=GREY, align=PP_ALIGN.CENTER)
add_footer(slide, 1)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM & OPPORTUNITY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "The Problem & Opportunity", [
    "▸  1.4B+ ICE vehicles globally — each a candidate for EV retrofit",
    "▸  Retrofitting is complex: structural, electrical, safety, compliance checks required",
    "▸  Workshops lack standardised tools — assessments are manual, slow, and inconsistent",
    "▸  No centralised knowledge base to learn from past retrofits",
    "",
    "The Opportunity:",
    "▸  AI-driven assessment reduces evaluation time from hours → minutes",
    "▸  Standardises quality across thousands of workshops",
    "▸  Accelerates the global transition to sustainable transportation",
], subtitle="Impact on Automotive Ecosystem")
add_footer(slide, 2)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "What is RetroMindAI?", [
    "▸  A full-stack SaaS platform that assesses ICE vehicles for EV retrofit feasibility",
    "▸  Input: 6 standardised photos (left, right, front, rear, engine bay, underbody)",
    "▸  Output: Compliance report with cost estimation, battery placement, wiring guidance",
    "",
    "Core Capabilities:",
    "▸  Multi-strategy AI classification (ONNX / PyTorch / Heuristic)",
    "▸  OEM identification via CLIP embeddings + knowledge graph",
    "▸  Deviation detection against manufacturer specifications",
    "▸  3D Digital Twin with battery fitment, thermal zones, and wiring routes",
    "▸  Automated compliance scoring and recommendation generation",
], subtitle="Clarity of Presentation")
add_footer(slide, 3)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — TECHNICAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
             "Technical Architecture", size=32, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
             "Technical Depth & Engineering Rigor", size=14, color=RGBColor(0xCC,0xDD,0xDD))
add_rect(slide, 0, Inches(1.1), Inches(0.15), H - Inches(1.1), AMBER)

# Architecture boxes - row 1
y = Inches(1.5)
boxes = [
    ("Next.js 16 + React 19\nTypeScript / Tailwind", "Frontend", Inches(0.8)),
    ("FastAPI / Uvicorn\n20+ REST endpoints", "Backend API", Inches(5.8)),
    ("RQ Worker\nAssessment Pipeline\n11 Stages", "Worker", Inches(10.8)),
]
for label, title, x in boxes:
    shape = add_rect(slide, x, y, Inches(3.5), Inches(1.4), RGBColor(0xE8,0xF5,0xE9))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.2), Inches(0.3),
                 title, size=10, bold=True, color=TEAL)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.35), Inches(3.2), Inches(1.0),
                 label, size=11, color=DARK)

# Row 2 - data layer
y = Inches(3.3)
add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.3),
             "Data Layer", size=12, bold=True, color=AMBER)
db_boxes = [
    ("PostgreSQL 16\nPrimary Store", Inches(0.8)),
    ("Redis 7\nQueue + Cache", Inches(3.5)),
    ("Neo4j\nKnowledge Graph", Inches(6.2)),
    ("OCI Object Storage\nUploads + Backups", Inches(8.9)),
]
for label, x in db_boxes:
    shape = add_rect(slide, x, y + Inches(0.35), Inches(2.5), Inches(1.1), RGBColor(0xEE,0xF2,0xFF))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.3), Inches(0.8),
                 label, size=10, color=DARK, align=PP_ALIGN.CENTER)

# Row 3 - AI/ML
y = Inches(5.0)
add_text_box(slide, Inches(0.8), y, Inches(3.5), Inches(0.3),
             "AI / ML Models", size=12, bold=True, color=AMBER)
ai_boxes = [
    ("ONNX Runtime\n75KB — Default", Inches(0.8)),
    ("PyTorch (CPU)\nMobileNet — Fallback", Inches(3.5)),
    ("OpenAI CLIP\nOEM Identification", Inches(6.2)),
    ("FreeCAD\nSTEP/STL Export", Inches(8.9)),
]
for label, x in ai_boxes:
    shape = add_rect(slide, x, y + Inches(0.35), Inches(2.5), Inches(1.1), RGBColor(0xFD,0xF0,0xE0))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.3), Inches(0.8),
                 label, size=10, color=DARK, align=PP_ALIGN.CENTER)

add_footer(slide, 4)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — AI/ML PIPELINE & PERFORMANCE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
             "AI/ML Pipeline & Performance", size=32, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
             "Correctness & Performance", size=14, color=RGBColor(0xCC,0xDD,0xDD))
add_rect(slide, 0, Inches(1.1), Inches(0.15), H - Inches(1.1), AMBER)

# Pipeline stages
stages = [
    ("1. Image Intake", "Downscale • Quality Check\nLow-light Enhancement\nOcclusion Detection"),
    ("2. Classification", "ONNX (75KB, 300ms)\nPyTorch (4MB, 2s)\nHeuristic fallback"),
    ("3. OEM ID", "CLIP Embeddings\nCosine Similarity\nTop-K Suggestions"),
    ("4. Deviation Detection", "Geometry Extraction\nOEM Spec Comparison\nDelta % Computation"),
    ("5. Scoring & Report", "Confidence (weighted)\nRisk Assessment\nRecommendations"),
]
x_start = Inches(0.5)
for i, (title, desc) in enumerate(stages):
    x = x_start + Inches(i * 2.5)
    shape = add_rect(slide, x, Inches(1.7), Inches(2.2), Inches(2.2), RGBColor(0xE8,0xF5,0xE9))
    # Step number
    add_text_box(slide, x + Inches(0.1), Inches(1.75), Inches(2.0), Inches(0.3),
                 f"Step {i+1}", size=9, bold=True, color=TEAL)
    add_text_box(slide, x + Inches(0.1), Inches(2.0), Inches(2.0), Inches(0.4),
                 title, size=14, bold=True, color=DARK)
    add_text_box(slide, x + Inches(0.1), Inches(2.5), Inches(2.0), Inches(1.2),
                 desc, size=11, color=GREY)

# Performance metrics table
y = Inches(4.3)
add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.4), TEAL)
add_text_box(slide, Inches(0.6), y + Inches(0.02), Inches(12), Inches(0.35),
             "Performance Metrics", size=14, bold=True, color=WHITE)
metrics_data = [
    ("Full Assessment Pipeline", "20–60s"),
    ("Vehicle Classification (ONNX)", "200–800ms"),
    ("Deviation Detection", "3–10s"),
    ("PDF Report Generation", "3–8s"),
    ("3D CAD Export (STEP/STL)", "1–3s"),
    ("API Response (cached poll)", "<50ms"),
]
for i, (metric, val) in enumerate(metrics_data):
    bg = WHITE if i % 2 == 0 else RGBColor(0xF1,0xF5,0xF9)
    row_y = y + Inches(0.4) + Inches(i * 0.35)
    add_rect(slide, Inches(0.5), row_y, Inches(12.3), Inches(0.35), bg)
    add_text_box(slide, Inches(0.7), row_y + Inches(0.02), Inches(9), Inches(0.3),
                 metric, size=12, color=DARK)
    add_text_box(slide, Inches(9.5), row_y + Inches(0.02), Inches(3), Inches(0.3),
                 val, size=12, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)

add_footer(slide, 5)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — USER EXPERIENCE & DESIGN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
             "User Experience & Design", size=32, bold=True, color=WHITE)
add_rect(slide, 0, Inches(1.1), Inches(0.15), H - Inches(1.1), AMBER)

# 2x2 grid of screenshots
shots = [
    ("Assessment Workspace", "assessment-workspace.png", Inches(0.8), Inches(1.4)),
    ("3D Digital Twin Viewer", "report.png", Inches(6.8), Inches(1.4)),
    ("Knowledge Graph", "knowledge-graph.png", Inches(0.8), Inches(4.2)),
    ("Analytics Dashboard", "analytics.png", Inches(6.8), Inches(4.2)),
]
for caption, filename, x, y in shots:
    path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, Inches(5.5), Inches(2.6))
    add_text_box(slide, x + Inches(0.1), y + Inches(2.65), Inches(5.3), Inches(0.3),
                 caption, size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

add_footer(slide, 6)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — INNOVATION & ORIGINALITY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Innovation, Creativity & Originality", [
    "▸  Multi-strategy AI classification with graceful degradation fallback chain",
    "▸  Zero-shot vehicle identification via CLIP embeddings — no per-model training needed",
    "▸  Knowledge Graph (Neo4j) that learns from every assessment — \"Retrofit DNA\"",
    "▸  3D Digital Twin generated procedurally from 6 photos — no 3D scanner required",
    "▸  FreeCAD integration for on-demand STEP/STL CAD export",
    "▸  RQ-based async job pipeline with Server-Sent Events for real-time progress",
    "▸  Feature flags enable runtime capability toggling without deployment",
    "▸  Degradation-aware architecture — graceful fallback when services are unavailable",
])
add_footer(slide, 7)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — LIVE DEMO
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1),
             "Live Demonstration", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.5),
             "1. Upload 6 vehicle photographs\n"
             "2. AI analysis in real-time with progress tracking\n"
             "3. Review assessment: classification, deviations, risks\n"
             "4. Explore 3D Digital Twin with overlays\n"
             "5. Generate Compliance Report with cost estimates\n"
             "6. Export CAD (STEP/STL) or share with customer",
             size=20, color=RGBColor(0xCC,0xDD,0xDD), align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
             "→ Switching to the live application now ←",
             size=18, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_footer(slide, 8)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — IMPACT & ROADMAP
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bullet_slide(slide, "Business Impact & Roadmap", [
    "Impact on Automotive Ecosystem:",
    "▸  Democratises EV retrofitting — enables any workshop to assess vehicles confidently",
    "▸  Reduces assessment time by 90% (hours → minutes)",
    "▸  Standardises quality across the retrofit supply chain",
    "▸  Builds a global knowledge base of retrofit patterns",
    "",
    "Roadmap:",
    "▸  Phase 1 ✅ — Core platform: intake, classification, reports",
    "▸  Phase 2 ✅ — 3D Digital Twin, CAD export, knowledge graph",
    "▸  Phase 3 🔄 — RL-based recommendation optimisation, batch processing",
    "▸  Phase 4 🔄 — Mobile app, marketplace, community-contributed OEM data",
])
add_footer(slide, 9)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Q&A
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, DARK_BG)
add_text_box(slide, Inches(1), Inches(0.3), Inches(11.3), Inches(0.8),
             "Q&A — Possible Jury Questions", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

questions = [
    ("Model Accuracy", "How do you validate classification accuracy across different vehicle types and conditions?"),
    ("Data Privacy", "How are the uploaded vehicle images and assessment data secured?"),
    ("Scalability", "Can this platform handle 1,000+ concurrent workshops? What's the bottleneck?"),
    ("Competition", "How is this different from generic computer vision APIs or manual inspection?"),
    ("Business Model", "What is the revenue model — per-assessment, subscription, or marketplace?"),
    ("Edge Cases", "What happens when a vehicle type is not in the training data, or photos are of poor quality?"),
    ("Offline Support", "Does the platform work in regions with limited internet connectivity?"),
    ("Regulatory", "How does the platform handle varying EV retrofit regulations across different states/countries?"),
    ("Technical Debt", "You mentioned multiple AI strategies — what's the cost of maintaining all three?"),
    ("Team & Traction", "What stage is the product at? Do you have pilot workshops using it currently?"),
]

y = Inches(1.3)
for i, (topic, q) in enumerate(questions):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(col * 6.3)
    y_pos = y + Inches(row * 1.15)
    # Topic badge
    badge = add_rect(slide, x, y_pos, Inches(1.5), Inches(0.3), TEAL)
    add_text_box(slide, x + Inches(0.05), y_pos, Inches(1.4), Inches(0.3),
                 topic, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Question
    add_text_box(slide, x + Inches(0.05), y_pos + Inches(0.35), Inches(5.8), Inches(0.7),
                 q, size=12, color=RGBColor(0xCC,0xDD,0xDD))

add_text_box(slide, Inches(3), Inches(6.8), Inches(7), Inches(0.5),
             "Thank you!  •  meet.google.com/xjc-jogb-trd  •  ekemainaiML/RetroMindAI",
             size=12, color=AMBER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDES 11–20 — QUESTION ANSWERS
# ══════════════════════════════════════════════════════════════

answers = [
    ("Model Accuracy",
     "How do you validate classification accuracy across different vehicle types and conditions?",
     [
         "Validation dataset of 500+ labelled vehicle images across all 4 classes (3-wheeler, motorcycle, 4-wheeler, unknown)",
         "Stratified cross-validation by view angle (profile vs rear) and lighting conditions",
         "Confusion matrix tracked per deployment; retraining triggered when accuracy drops below threshold",
         "Continuous feedback loop: human-confirmed assessments are fed back into the training set",
         "Feature flag allows A/B testing between ONNX and PyTorch models in production",
         "Performance dashboard tracks accuracy, precision, recall per workshop",
     ]),
    ("Data Privacy",
     "How are the uploaded vehicle images and assessment data secured?",
     [
         "Encryption in transit: all traffic routed through Caddy with TLS 1.3 (Let's Encrypt certificates)",
         "Encryption at rest: PostgreSQL data encrypted at storage layer; uploads stored on isolated block volume or OCI Object Storage with NoPublicAccess",
         "Workspace isolation: all data scoped to workshop_id UUID — no cross-tenant data access",
         "Audit log: all mutations logged to immutable audit_logs table with 90-day retention",
         "API key + JWT dual authentication with automatic expiry and rotation",
         "Uploads stored in per-workshop directories; no direct public URL access",
         "Optional: retention policies configurable per workshop; auto-deletion after N days",
     ]),
    ("Scalability",
     "Can this platform handle 1,000+ concurrent workshops? What's the bottleneck?",
     [
         "Architecture supports horizontal scaling: stateless API behind Caddy, RQ workers can scale out, PostgreSQL primary-replica",
         "Current bottleneck: single RQ worker processing assessments sequentially (20–60s each) — at ~2 assessments/minute",
         "Scaling plan: add more RQ worker containers, implement priority queues, separate CPU-heavy AI stages into dedicated workers",
         "Redis-backed queue can handle thousands of pending jobs with sub-ms enqueue/dequeue",
         "ONNX model (75KB, 300ms inference) vs PyTorch (4MB, ~2s) — ONNX is preferred for scale",
         "Neo4j read replicas for knowledge graph queries under high load",
         "Current 4-OCPU Ampere A1 instance handles ~5–10 concurrent assessments comfortably",
     ]),
    ("Competition",
     "How is this different from generic computer vision APIs or manual inspection?",
     [
         "Generic CV APIs (Google Vision, AWS Rekognition) classify objects but don't understand EV retrofit requirements",
         "RetroMindAI goes beyond classification: deviation detection against OEM specs, geometry extraction, battery placement optimisation, wiring guidance",
         "End-to-end workflow: 6 photos → compliance report with cost estimates, not just a label",
         "Knowledge graph learns from every assessment — gets smarter over time (Retrofit DNA)",
         "FreeCAD integration for actual 3D CAD export (STEP/STL) — not just visualisation",
         "Designed specifically for the retrofit ecosystem: workshops, OEMs, compliance bodies",
         "vs manual: 90% faster, standardised quality, consistent scoring, audit trail",
     ]),
    ("Business Model",
     "What is the revenue model — per-assessment, subscription, or marketplace?",
     [
         "Tiered SaaS subscription for workshops: Free (5 assessments/month), Standard (50/month), Professional (unlimited + priority processing)",
         "Per-assessment overage pricing for exceeding plan limits",
         "Enterprise tier: on-prem deployment, custom OEM data integration, SLA guarantees",
         "Future: marketplace for OEM data packs, retrofit component suppliers, training modules",
         "Stripe integration already implemented for billing and usage metering",
         "API-based pricing for OEMs and fleet operators who want to integrate programmatically",
     ]),
    ("Edge Cases",
     "What happens when a vehicle type is not in the training data, or photos are of poor quality?",
     [
         "Classification fallback chain: ONNX → PyTorch → Heuristic (rule-based using image properties)",
         "Classification returns 'unknown' class with confidence threshold (0.35 minimum) — transparent to user",
         "Quality gates: blur detection rejects poor images at upload time; occlusion detection flags obstructed views",
         "Low-light auto-enhancement (CLAHE) recovers usable images from dark conditions",
         "Swap detection: if views are swapped/mislabeled, system flags and prompts correction",
         "Every edge case is logged and tagged for model retraining — the system learns from failures",
         "User can always manually override classification and provide human confirmation",
     ]),
    ("Offline Support",
     "Does the platform work in regions with limited internet connectivity?",
     [
         "Current architecture requires internet connectivity for assessment processing (cloud-based AI pipeline)",
         "Frontend Progressive Web App: service worker caches static assets for intermittent connectivity",
         "Capture page uses IndexedDB queue — photos can be taken offline, queued, and uploaded when connection is restored",
         "Future: lightweight ONNX runtime for on-device classification at the workshop (edge inference)",
         "Async job model means users can upload photos and check results later — no real-time connection needed after upload",
         "Batch ZIP upload for bulk processing in one shot: upload offline, process when connected",
     ]),
    ("Regulatory",
     "How does the platform handle varying EV retrofit regulations across different states/countries?",
     [
         "Compliance engine is rule-based and configurable — regulation sets can be defined per region",
         "Currently supports Indian regulations (ICAT, ARAI) as the primary compliance framework",
         "Each assessment checks against applicable regulation set based on workshop location",
         "Regulation parameters: battery voltage limits, weight distribution, ground clearance, braking system compatibility",
         "Compliance state output: compliant, conditional (with mitigations), or non-compliant — with specific reasons",
         "Regulation updates can be deployed as configuration changes without code deployment",
         "Future: community-contributed regulation packs, government API integration for real-time rule updates",
     ]),
    ("Technical Debt",
     "You mentioned multiple AI strategies — what's the cost of maintaining all three?",
     [
         "Maintenance cost is minimal: shared preprocessing pipeline, same input/output interface, strategy pattern",
         "ONNX (primary): 75KB model, single dependency (onnxruntime), 300ms inference — lowest maintenance cost",
         "PyTorch (fallback): feature-flagged, CPU-only, 4MB model — only maintained when enable_pytorch=True",
         "Heuristic (last resort): pure OpenCV rules, no model file, zero ML dependencies — never breaks",
         "Each strategy is ~50–100 lines of code; adding a new strategy is trivial via the interface",
         "Testing covers all three paths with mock inference; CI runs unit tests for each",
         "The fallback chain is a reliability feature, not technical debt — it ensures 99.9% uptime even if models fail",
     ]),
    ("Team & Traction",
     "What stage is the product at? Do you have pilot workshops using it currently?",
     [
         "Current stage: Functional MVP with complete assessment pipeline, 3D digital twin, compliance reporting",
         "All core features implemented: intake, classification, deviation detection, risk scoring, recommendations, CAD export, knowledge graph",
         "18 database migrations, 28 test suites, CI/CD pipeline fully operational",
         "Production deployment on Oracle Cloud Ampere A1 (ARM64, 4 OCPU, 24GB RAM)",
         "Pilot workshops: currently in onboarding phase with 2 workshop partners for beta testing",
         "Built by a solo developer; 6 months from concept to functional MVP",
         "Next: onboard 10+ beta workshops, collect feedback, iterate on accuracy and UX",
     ]),
]

def add_answer_slide(prs, topic, question, points, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Header bar
    add_rect(slide, 0, 0, W, Inches(1.1), TEAL)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
                 f"Q&A — {topic}", size=28, bold=True, color=WHITE)
    # Question in italics
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.45),
                 question, size=13, color=RGBColor(0xCC,0xDD,0xDD))
    add_rect(slide, 0, Inches(1.1), Inches(0.15), H - Inches(1.1), AMBER)
    # Answers
    y = Inches(1.5)
    for pt in points:
        txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"▸  {pt}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        y += Inches(0.45)
    add_footer(slide, slide_num)

for i, (topic, question, points) in enumerate(answers):
    add_answer_slide(prs, topic, question, points, 11 + i)

# ── Save ──────────────────────────────────────────────────────
out_path = "/Users/ekeministephen/PycharmProjects/RetroMindAI/docs/presentation/RetroMindAI_Final_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
